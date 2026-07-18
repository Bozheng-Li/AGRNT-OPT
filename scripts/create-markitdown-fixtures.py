"""Create deterministic, local-only fixtures for MarkItDown integration tests."""

from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation


def write_pdf(target: Path) -> None:
    stream = b"BT /F1 18 Tf 72 720 Td (PDF fixture evidence) Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(stream), stream),
    ]
    body = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for number, value in enumerate(objects, start=1):
        offsets.append(len(body))
        body.extend(f"{number} 0 obj\n".encode())
        body.extend(value)
        body.extend(b"\nendobj\n")
    xref = len(body)
    body.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        body.extend(f"{offset:010d} 00000 n \n".encode())
    body.extend(f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode())
    target.write_bytes(body)


def write_docx(target: Path) -> None:
    content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""
    relationships = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""
    document = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body>
  <w:p><w:pPr><w:pStyle w:val="Title"/></w:pPr><w:r><w:t>DOCX fixture evidence</w:t></w:r></w:p>
  <w:p><w:r><w:t>Office container conversion passed.</w:t></w:r></w:p>
  <w:sectPr><w:pgSz w:w="12240" w:h="15840"/></w:sectPr>
</w:body></w:document>"""
    with zipfile.ZipFile(target, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", relationships)
        archive.writestr("word/document.xml", document)


def write_unsafe_docx_fixtures(output: Path) -> None:
    with zipfile.ZipFile(output / "traversal.docx", "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "types")
        archive.writestr("word/document.xml", "document")
        archive.writestr("../escape.txt", "escape")
    with zipfile.ZipFile(output / "bomb.docx", "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "types")
        archive.writestr("word/document.xml", b"x" * (33 * 1024 * 1024))
    with zipfile.ZipFile(output / "wrong-type.docx", "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "types")
        archive.writestr("xl/workbook.xml", "workbook")


def write_pptx(target: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPTX fixture evidence"
    slide.placeholders[1].text = "Presentation conversion passed."
    presentation.save(target)


def write_xlsx(target: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    sheet.append(["capability", "status"])
    sheet.append(["XLSX fixture evidence", "passed"])
    workbook.save(target)


def main() -> None:
    if len(sys.argv) != 2:
        raise SystemExit("usage: create-markitdown-fixtures.py OUTPUT_DIRECTORY")
    output = Path(sys.argv[1]).resolve()
    output.mkdir(parents=True, exist_ok=True)
    write_pdf(output / "evidence.pdf")
    write_docx(output / "evidence.docx")
    write_unsafe_docx_fixtures(output)
    write_pptx(output / "evidence.pptx")
    write_xlsx(output / "evidence.xlsx")
    (output / "evidence.json").write_text(json.dumps({"title": "JSON fixture evidence", "status": "passed"}), encoding="utf-8")
    (output / "evidence.htm").write_text("<h1>HTM fixture evidence</h1><p>Alias extension passed.</p>", encoding="utf-8")
    (output / "evidence.md").write_text("# Markdown fixture evidence\n\nSource preservation passed.\n", encoding="utf-8")
    (output / "evidence.txt").write_text("Plain text fixture evidence\n", encoding="utf-8")


if __name__ == "__main__":
    main()
